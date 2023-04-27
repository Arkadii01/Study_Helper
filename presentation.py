import collections
import collections.abc  # Для предотвращения неизвестной ошибки
from pptx import Presentation
import docx
import webbrowser
import os


def text_for_slides(docx_file):
    doc = docx.Document(docx_file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)    
    text = (' '.join(full_text)).split(' ')
    total_text = []
    slide_text = []
    links_text = []
    for word in text:
        if 'Источники:' not in word:
            slide_text.append(word)
            if '.' in word:
                if len(slide_text) >= 40:
                    total_text.append(slide_text)
                    slide_text = []
                    continue
        else:
            links_text.append(word[10:])
    if slide_text:
        total_text.append(slide_text)
    return total_text, links_text

docx_file = f'{input("Введите docx файл без .docx")}.docx'
total_text, links = text_for_slides(docx_file)

prs = Presentation()
need = input()

slide = prs.slides.add_slide(prs.slide_layouts[0])   # slide
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = need
subtitle.text = 'Работу выполнил:\n'

slide = prs.slides.add_slide(prs.slide_layouts[1]) # slide
shapes = slide.shapes
title_shape = shapes.title
text = slide.placeholders[1]
title_shape.text = 'Цели и задачи'
text.text = 'Основная цель - .\nЗадачи - .'

for i in range(len(total_text)):
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # slides
    text = slide.placeholders[1]
    text.text = ' '.join(total_text[i])

slide = prs.slides.add_slide(prs.slide_layouts[1])  # slide
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = f'Источники:\n{links}'

prs.save(f'{need}.pptx')

#img-s
webbrowser.open(f'https://yandex.ru/images/search?text={"%20".join(need.split(" "))}%20рисунки%20без%20текста&itype=png&from=tabbar')